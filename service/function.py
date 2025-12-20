import re
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_empty_slide(prs):
    """빈페이지 넣기, 간주하거나 곡 끝나고 분리용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)


def add_title_slide(prs, title):
    """제목 페이지 일반 가사보다 조금 아래로 되어 있음."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    box = slide.shapes.add_textbox(
        (prs.slide_width - Pt(1400)) / 2,
        (prs.slide_height - Pt(200)) / 2,
        Pt(1400),
        Pt(200),
    )

    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = title

    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def parse_song_form(song_form: str):
    """
    송폼을 리스트 형식으로 바꿔주는 파싱 함수
    """
    return re.findall(r"\([^)]+\)|[A-Z][0-9]*", song_form)


def add_lyrics_slide(prs, slide_text: str):
    """
    가사 슬라이드 넣기 // 추후에 뭐 폰트나 색 이런거 여기서 수정
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    BOX_WIDTH = Pt(1400)
    BOX_HEIGHT = Pt(300)
    LEFT = (prs.slide_width - BOX_WIDTH) / 2
    TOP = Pt(0)

    box = slide.shapes.add_textbox(LEFT, TOP, BOX_WIDTH, BOX_HEIGHT)

    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = [l.strip() for l in slide_text.split("\n") if l.strip()]

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def render_part(prs, part_text: str):
    """
    파트 문자열 전체를 받아서
    // 기준으로 슬라이드 생성
    """
    slides = [s.strip() for s in part_text.split("//") if s.strip()]
    for slide_text in slides:
        add_lyrics_slide(prs, slide_text)


def render_song(prs, song: dict):
    """
    곡 1개 전체 렌더링
    """
    parts = song["parts"]
    song_form = song["song_form"]
    title = song.get("title")

    if title:
        add_title_slide(prs, title)

    tokens = parse_song_form(song_form)

    for token in tokens:
        if token.startswith("("):
            add_empty_slide(prs)
            continue

        text = parts.get(token, "").strip()
        if not text:
            continue

        render_part(prs, text)

    # 곡 끝 여백
    add_empty_slide(prs)
    add_empty_slide(prs)
    add_empty_slide(prs)
